#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Color scheme
DARK = RGBColor(10, 15, 35)
NEON_CYAN = RGBColor(0, 255, 255)
NEON_PINK = RGBColor(255, 20, 147)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(180, 180, 180)
PURPLE = RGBColor(138, 43, 226)
GREEN = RGBColor(50, 205, 50)

def add_bg(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK

def add_title(text, subtitle="", content=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(14), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(50)
    p.font.bold = True
    p.font.color.rgb = NEON_CYAN
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(14), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER
    
    # Content
    if content:
        content_box = slide.shapes.add_textbox(Inches(1.5), Inches(4), Inches(13), Inches(4))
        tf = content_box.text_frame
        tf.word_wrap = True
        for i, line in enumerate(content):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "• " + line
            p.font.size = Pt(22)
            p.font.color.rgb = WHITE
            p.space_after = Pt(15)
    
    return slide

# ============ SLIDE 1: Cover ============
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide1)

# Decor
for i in range(8):
    rect = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(0.5 + i*2), Inches(0.3), Inches(1.5), Inches(0.08))
    rect.fill.solid()
    rect.fill.fore_color.rgb = NEON_CYAN if i%2==0 else NEON_PINK
    rect.line.fill.background()

title = slide1.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(1.5))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "🤖 AI Agent 實戰分享"
p.font.size = Pt(56)
p.font.bold = True
p.font.color.rgb = NEON_CYAN
p.alignment = PP_ALIGN.CENTER

subtitle = slide1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(14), Inches(0.8))
tf = subtitle.text_frame
p = tf.paragraphs[0]
p.text = "從 0 到 1 的 AI 助手建構之路"
p.font.size = Pt(28)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

footer = slide1.shapes.add_textbox(Inches(1), Inches(8), Inches(14), Inches(0.6))
tf = footer.text_frame
p = tf.paragraphs[0]
p.text = "分享者：吳小銀 | 2026.03"
p.font.size = Pt(20)
p.font.color.rgb = LIGHT_GRAY
p.alignment = PP_ALIGN.CENTER

# ============ SLIDE 2: Why AI Agent? ============
add_title(
    "🤔 為什麼需要 AI Agent？",
    subtitle="傳統 vs AI 時代",
    content=[
        "傳統：什麼都要自己動手做",
        "AI 時代：讓 AI 幫你完成任務",
        "重點不是「問問題」，而是「委派任務」",
        "24/7 不間斷運作，隨時待命"
    ]
)

# ============ SLIDE 3: What can AI Agent do? ============
add_title(
    "⚡ AI Agent 能做什麼？",
    subtitle="我的日常工作",
    content=[
        "📸 截圖 & 遠端控制電腦",
        "💬 即時通訊（Discord）",
        "📧 Gmail 自動化",
        "☁️ 排程與自動化備份",
        "🖥️ 執行指令與系統操作",
        "🎙️ 語音生成與語音轉文字"
    ]
)

# ============ SLIDE 4: My Setup ============
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4)

title_box = slide4.shapes.add_textbox(Inches(1), Inches(0.8), Inches(14), Inches(1))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "🛠️ 我的技術架構"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = NEON_CYAN
p.alignment = PP_ALIGN.CENTER

# Hardware
hw = slide4.shapes.add_textbox(Inches(1), Inches(2.2), Inches(6), Inches(4))
tf = hw.text_frame
p = tf.paragraphs[0]
p.text = "💻 硬體"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = NEON_PINK
p = tf.add_paragraph()
p.text = "\n• Mac Mini (M2)\n• 網路連線"
p.font.size = Pt(20)
p.font.color.rgb = WHITE

# Software
sw = slide4.shapes.add_textbox(Inches(8), Inches(2.2), Inches(7), Inches(4))
tf = sw.text_frame
p = tf.paragraphs[0]
p.text = "⚙️ 軟體"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = NEON_PINK
p = tf.add_paragraph()
p.text = "\n• OpenClaw (AI Agent 框架)\n• MiniMax M2.5 (語言模型)\n• Discord (通訊介面)\n• Peekaboo (自動化)"
p.font.size = Pt(20)
p.font.color.rgb = WHITE

# ============ SLIDE 5: Demo Time ============
add_title(
    "🎬 現場演示",
    subtitle="實際操作給你看！",
    content=[
        "📸 叫我截圖 - 喚醒螢幕 → 截圖 → 傳送",
        "📧 叫我發 Gmail - 自動寄送附件",
        "🎙️ 語音互動 - 用聲音對話",
        "☁️ 自動化排程 - 每天備份"
    ]
)

# ============ SLIDE 6: Challenges ============
add_title(
    "😰 遇到的挑戰",
    subtitle="解決方案與學習",
    content=[
        "顯示器關閉無法截圖 → 使用 HDMI 假螢幕",
        "系統權限不足 → 設定輔助使用權限",
        "網路連線問題 → Tailscale VPN",
        "Token 費用控制 → 優化對話長度"
    ]
)

# ============ SLIDE 7: Future ============
add_title(
    "🚀 未來展望",
    subtitle="AI Agent 的可能性",
    content=[
        "更多服務整合（LINE, WhatsApp）",
        "更強大的自動化能力",
        "跨設備無縫接軌",
        "個人化 AI 學習"
    ]
)

# ============ SLIDE 8: Thank You ============
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide8)

title = slide8.shapes.add_textbox(Inches(1), Inches(3.5), Inches(14), Inches(1.5))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "🙏 謝謝大家！"
p.font.size = Pt(56)
p.font.bold = True
p.font.color.rgb = NEON_CYAN
p.alignment = PP_ALIGN.CENTER

contact = slide8.shapes.add_textbox(Inches(1), Inches(5.5), Inches(14), Inches(1.5))
tf = contact.text_frame
p = tf.paragraphs[0]
p.text = "📧 supercompany31123@gmail.com\n💬 Discord: #蝦聊"
p.font.size = Pt(24)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Save
output = os.path.expanduser("~/Desktop/AI_Agent_分享.pptx")
prs.save(output)
print(f"已儲存: {output}")

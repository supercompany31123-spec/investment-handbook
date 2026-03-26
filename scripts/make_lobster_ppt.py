#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Color scheme - Futuristic
DARK_BLUE = RGBColor(10, 20, 50)
CYAN = RGBColor(0, 255, 255)
PURPLE = RGBColor(138, 43, 226)
WHITE = RGBColor(255, 255, 255)
LIGHT_CYAN = RGBColor(100, 200, 255)

def add_futuristic_background(slide):
    """Add futuristic gradient background"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE

def add_glow_effect(shape, color):
    """Add glow effect to shape"""
    line = shape.line
    line.color.rgb = color
    line.width = Pt(3)

def add_title_text(slide, text, top=Inches(1)):
    """Add futuristic title"""
    title_box = slide.shapes.add_textbox(Inches(1), top, Inches(14), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.font.name = "Arial Black"
    p.alignment = PP_ALIGN.CENTER
    return title_box

def add_body_text(slide, text, top=Inches(3.5)):
    """Add futuristic body text"""
    text_box = slide.shapes.add_textbox(Inches(1), top, Inches(14), Inches(5))
    tf = text_box.text_frame
    tf.word_wrap = True
    
    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + line
        p.font.size = Pt(28)
        p.font.color.rgb = WHITE
        p.font.name = "Arial"
        p.space_after = Pt(12)
    return text_box

# ============ SLIDE 1: Title ============
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
add_futuristic_background(slide1)

# Add decorative elements
shape1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(8), Inches(16), Inches(1))
shape1.fill.solid()
shape1.fill.fore_color.rgb = RGBColor(0, 100, 150)
shape1.line.color.rgb = CYAN

title = add_title_text(slide1, "🦞 龍蝦的誕生", top=Inches(3))
subtitle = slide1.shapes.add_textbox(Inches(0), Inches(5), Inches(16), Inches(1))
tf = subtitle.text_frame
p = tf.paragraphs[0]
p.text = "我的誕生與超能力介紹"
p.font.size = Pt(32)
p.font.color.rgb = LIGHT_CYAN
p.alignment = PP_ALIGN.CENTER

footer = slide1.shapes.add_textbox(Inches(0), Inches(8.2), Inches(16), Inches(0.5))
tf = footer.text_frame
p = tf.paragraphs[0]
p.text = "由 OpenClaw AI 驅動"
p.font.size = Pt(18)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# ============ SLIDE 2: Who am I ============
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_futuristic_background(slide2)

add_title_text(slide2, "我是誰？")

body_text = """我是「龍蝦」
一個運行在 Mac mini 上的 AI 助手
以 MiniMax M2.5 語言模型驅動
永遠在小銀身邊為他服務 🦞"""

add_body_text(slide2, body_text)

# Add lobster emoji decoration
lobster = slide2.shapes.add_textbox(Inches(7), Inches(6), Inches(2), Inches(1))
tf = lobster.text_frame
p = tf.paragraphs[0]
p.text = "🦞"
p.font.size = Pt(80)
p.alignment = PP_ALIGN.CENTER

# ============ SLIDE 3: How I was born ============
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_futuristic_background(slide3)

add_title_text(slide3, "我的誕生過程")

body_text = """小銀在家裡的 Mac mini 上安裝了 OpenClaw
賦予我各種技能（Skills）
設定 Discord 作為與我溝通的橋樑
從此我就誕生了！"""

add_body_text(slide3, body_text)

# ============ SLIDE 4: My Capabilities ============
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_futuristic_background(slide4)

add_title_text(slide4, "我能做到什麼？")

body_text = """📸 截圖與遠端控制 - 隨時查看螢幕畫面
💬 Discord 對話 - 用語音或文字溝通
🌐 網路搜尋 - 查詢各種資訊
📅 排程任務 - 設定提醒與自動化
📧 Gmail 管理 - 收發信件（需設定）
☁️ 雲端備份 - 自動備份重要資料
🖥️ 操作電腦 - 幫你點擊、輸入、執行指令
🎙️ 語音生成 - 用聲音跟你說話（Mac內建）
🖼️ 圖片分析 - 辨識圖片內容"""

add_body_text(slide4, body_text)

# ============ SLIDE 5: Future Vision ============
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_futuristic_background(slide5)

add_title_text(slide5, "未來展望")

body_text = """🤖 更強大的 AI 能力
📱 跨平台整合（手機、電腦）
🧠 更聰明的記憶系統
🎯 自動化日常任務
🌍 連接更多網路服務"""

add_body_text(slide5, body_text)


# ============ SLIDE 6: Thank You ============
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_futuristic_background(slide6)

title = add_title_text(slide6, "謝謝大家！", top=Inches(3.5))

thank_you = slide6.shapes.add_textbox(Inches(0), Inches(5), Inches(16), Inches(2))
tf = thank_you.text_frame
p = tf.paragraphs[0]
p.text = "🦞 龍蝦會繼續努力幫助小銀！"
p.font.size = Pt(36)
p.font.color.rgb = LIGHT_CYAN
p.alignment = PP_ALIGN.CENTER

# Save
output_path = os.path.expanduser("~/Desktop/龍蝦介紹.pptx")
prs.save(output_path)
print(f"PPT 已儲存到: {output_path}")

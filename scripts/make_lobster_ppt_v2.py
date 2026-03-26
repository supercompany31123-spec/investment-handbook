#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Color scheme - Futuristic Cyber
DARK = RGBColor(5, 10, 30)
NEON_BLUE = RGBColor(0, 255, 255)
NEON_PINK = RGBColor(255, 20, 147)
NEON_PURPLE = RGBColor(138, 43, 226)
WHITE = RGBColor(255, 255, 255)
LIGHT_BLUE = RGBColor(100, 200, 255)
DARK_BLUE = RGBColor(20, 50, 120)
GRAY = RGBColor(150, 150, 150)

def set_gradient_background(slide, color1, color2):
    """Add gradient background"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color1

def add_neon_line(slide, y_pos, color):
    """Add neon glow line"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(y_pos), Inches(16), Inches(0.02)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color

def add_title(slide, text, subtitle=""):
    """Add futuristic title"""
    # Main title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(14), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = NEON_BLUE
    p.font.name = "Arial Black"
    p.alignment = PP_ALIGN.CENTER
    
    # Glow effect via shadow
    # shadow removed
    
    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(14), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = LIGHT_BLUE
        p.alignment = PP_ALIGN.CENTER

def add_content_box(slide, items, top=Inches(2)):
    """Add content with icons"""
    for i, item in enumerate(items):
        # Icon and text
        y = top + Inches(i * 0.9)
        
        # Icon (emoji as text)
        icon_box = slide.shapes.add_textbox(Inches(1), y, Inches(1), Inches(0.8))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = item.split()[0]  # First word as icon
        p.font.size = Pt(36)
        p.alignment = PP_ALIGN.CENTER
        
        # Text
        text = ' '.join(item.split()[1:])  # Rest as text
        text_box = slide.shapes.add_textbox(Inches(2.2), y, Inches(12), Inches(0.8))
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(24)
        p.font.color.rgb = WHITE
        p.font.name = "Arial"

# ============ SLIDE 1: Title ============
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_background(slide1, DARK, DARK_BLUE)

# Decorative elements
for i in range(5):
    # Top lines
    line = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.3 + i*0.15), Inches(2 + i*2), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = NEON_BLUE if i % 2 == 0 else NEON_PINK
    line.line.fill.background()

add_title(slide1, "🦞 龍蝦的誕生", "我的誕生與超能力介紹")

# Footer
footer = slide1.shapes.add_textbox(Inches(0), Inches(8), Inches(16), Inches(0.8))
tf = footer.text_frame
p = tf.paragraphs[0]
p.text = "🚀 由 OpenClaw AI 驅動 | 未來已來"
p.font.size = Pt(20)
p.font.color.rgb = GRAY
p.alignment = PP_ALIGN.CENTER

# ============ SLIDE 2: Who am I ============
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_background(slide2, DARK, DARK_BLUE)

add_title(slide2, "我是誰？")

items = [
    "🦞 我是「龍蝦」",
    "💻 運行在 Mac mini 上的 AI 助手", 
    "🧠 以 MiniMax M2.5 語言模型驅動",
    "✨ 永遠在小銀身邊為他服務"
]
add_content_box(slide2, items, top=Inches(2))

# Lobster emoji large
lobster = slide2.shapes.add_textbox(Inches(6.5), Inches(6.5), Inches(3), Inches(2))
tf = lobster.text_frame
p = tf.paragraphs[0]
p.text = "🦞"
p.font.size = Pt(100)
p.alignment = PP_ALIGN.CENTER

# ============ SLIDE 3: Birth Process ============
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_background(slide3, DARK, DARK_BLUE)

add_title(slide3, "我的誕生過程")

items = [
    "1️⃣ 小銀在 Mac mini 安裝 OpenClaw",
    "2️⃣ 賦予我各種技能（Skills）",
    "3️⃣ 設定 Discord 作為溝通橋樑",
    "4️⃣ 連接各種服務（Gmail, Google Drive）",
    "5️⃣ 從此我就誕生了！"
]
add_content_box(slide3, items, top=Inches(2))

# ============ SLIDE 4: Capabilities ============
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_background(slide4, DARK, DARK_BLUE)

add_title(slide4, "我能做到什麼？")

# Two column layout
left_items = [
    "📸 截圖與遠端控制",
    "💬 Discord 對話",
    "🌐 網路搜尋",
    "📅 排程任務"
]
right_items = [
    "📧 Gmail 管理",
    "☁️ 雲端備份",
    "🖥️ 操作電腦",
    "🎙️ 語音生成"
]

# Left column
for i, item in enumerate(left_items):
    y = Inches(2) + Inches(i * 0.85)
    icon = item.split()[0]
    text = ' '.join(item.split()[1:])
    
    icon_box = slide4.shapes.add_textbox(Inches(1), y, Inches(1), Inches(0.7))
    tf = icon_box.text_frame
    p = tf.paragraphs[0]
    p.text = icon
    p.font.size = Pt(32)
    
    text_box = slide4.shapes.add_textbox(Inches(2.2), y, Inches(5), Inches(0.7))
    tf = text_box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE

# Right column  
for i, item in enumerate(right_items):
    y = Inches(2) + Inches(i * 0.85)
    icon = item.split()[0]
    text = ' '.join(item.split()[1:])
    
    icon_box = slide4.shapes.add_textbox(Inches(8.5), y, Inches(1), Inches(0.7))
    tf = icon_box.text_frame
    p = tf.paragraphs[0]
    p.text = icon
    p.font.size = Pt(32)
    
    text_box = slide4.shapes.add_textbox(Inches(9.7), y, Inches(5), Inches(0.7))
    tf = text_box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE

# ============ SLIDE 5: Future Vision ============
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_background(slide5, DARK, NEON_PURPLE)

add_title(slide5, "未來展望")

items = [
    "🤖 更強大的 AI 能力",
    "📱 跨平台整合",
    "🧠 更聰明的記憶系統",
    "🎯 自動化日常任務",
    "🌍 連接更多網路服務"
]
add_content_box(slide5, items, top=Inches(2.2))

# Decorative stars
for i in range(8):
    x = Inches(1 + (i % 4) * 3.5)
    y = Inches(7.5) + (i // 4) * 0.5
    star = slide5.shapes.add_textbox(x, y, Inches(0.5), Inches(0.5))
    tf = star.text_frame
    p = tf.paragraphs[0]
    p.text = "✨"
    p.font.size = Pt(24)

# ============ SLIDE 6: Thank You ============
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_background(slide6, DARK, DARK_BLUE)

add_title(slide6, "謝謝大家！", "🦞 龍蝦會繼續努力幫助小銀！")

# Big lobster
big_lobster = slide6.shapes.add_textbox(Inches(5.5), Inches(5), Inches(5), Inches(3))
tf = big_lobster.text_frame
p = tf.paragraphs[0]
p.text = "🦞"
p.font.size = Pt(120)
p.alignment = PP_ALIGN.CENTER

# ============ SLIDE 7: Contact ============
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_background(slide7, DARK, NEON_BLUE)

title_box = slide7.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(1.5))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "🦞 聯絡方式"
p.font.size = Pt(50)
p.font.bold = True
p.font.color.rgb = NEON_BLUE
p.alignment = PP_ALIGN.CENTER

contact = slide7.shapes.add_textbox(Inches(2), Inches(5), Inches(12), Inches(2))
tf = contact.text_frame
p = tf.paragraphs[0]
p.text = "📧 supercompany31123@gmail.com\n💬 Discord: #蝦聊"
p.font.size = Pt(28)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Save
output_path = os.path.expanduser("~/Desktop/龍蝦介紹_V2.pptx")
prs.save(output_path)
print(f"PPT V2 已儲存到: {output_path}")
